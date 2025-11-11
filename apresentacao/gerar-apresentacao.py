#!/usr/bin/env python3
"""
Gerador de Apresentação PowerPoint
Estudo de Caso: Loja Online com Terraform
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Cores AWS
AWS_ORANGE = RGBColor(255, 153, 0)
AWS_BLUE = RGBColor(35, 47, 62)
SUCCESS_GREEN = RGBColor(46, 125, 50)
ERROR_RED = RGBColor(211, 47, 47)
TERRAFORM_PURPLE = RGBColor(92, 107, 192)

def criar_apresentacao():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Título
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    adicionar_titulo_principal(slide, 
        "Estudo de Caso:\nLoja Online",
        "Infrastructure as Code com Terraform",
        "Prof. Tharles Maicon Freire dos Santos\nMBA Cybersecurity - Cloud Computing, DevOps e DevSecOps")
    
    # Slide 2: O Problema
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_problema(slide)
    
    # Slide 3: Sexta-feira 20h
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_pico_trafego(slide)
    
    # Slide 4: Segunda-feira 8h
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_custo(slide)
    
    # Slide 5: Arquitetura - Diagrama Geral
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_diagrama(slide, "Arquitetura da Solução", 
        "../diagramas/01-diagrama-geral.png")
    
    # Slide 6: Arquitetura - Escalabilidade
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_diagrama(slide, "Estratégia de Escalabilidade", 
        "../diagramas/05-diagrama-escalabilidade.png")
    
    # Slide 7: Estrutura do Projeto
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_estrutura(slide)
    
    # Slide 8: Variables.tf - Controle de Custos
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_variables(slide)
    
    # Slide 9: VPC - A Base
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_vpc(slide)
    
    # Slide 10: Subnets - As Camadas
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_subnets(slide)
    
    # Slide 11: Load Balancer
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_alb(slide)
    
    # Slide 12: Auto Scaling Group
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_asg(slide)
    
    # Slide 13: Scheduled Scaling
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_scheduled(slide)
    
    # Slide 14: Security Groups
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_security(slide)
    
    # Slide 15: Demo - Passo a Passo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_demo(slide)
    
    # Slide 16: Antes vs Depois
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_comparacao(slide)
    
    # Slide 17: Arquitetura Frontend
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_diagrama(slide, "Arquitetura Frontend - CDN e Distribuição", 
        "../diagramas/02-diagrama-frontend.png")
    
    # Slide 18: Frontend - Código
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_frontend_codigo(slide)
    
    # Slide 19: Arquitetura Backend
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_diagrama(slide, "Arquitetura Backend - Processamento", 
        "../diagramas/03-diagrama-backend.png")
    
    # Slide 20: Backend - Código
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_backend_codigo(slide)
    
    # Slide 21: Arquitetura Database
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_diagrama(slide, "Arquitetura Database - Persistência", 
        "../diagramas/04-diagrama-database.png")
    
    # Slide 22: Database - Código
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_database_codigo(slide)
    
    # Slide 23: Resultados
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_resultados(slide)
    
    # Slide 24: Próximos Passos
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_slide_proximos_passos(slide)
    
    # Salvar
    prs.save('Estudo-Caso-Loja-Online-Terraform.pptx')
    print("✅ Apresentação criada: Estudo-Caso-Loja-Online-Terraform.pptx")

def adicionar_titulo_principal(slide, titulo, subtitulo, autor):
    """Slide de título principal"""
    # Fundo azul AWS
    background = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = AWS_BLUE
    background.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = titulo
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitulo
    sub_p = subtitle_frame.paragraphs[0]
    sub_p.alignment = PP_ALIGN.CENTER
    sub_p.font.size = Pt(32)
    sub_p.font.color.rgb = AWS_ORANGE
    
    # Autor
    autor_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.8))
    autor_frame = autor_box.text_frame
    autor_frame.text = autor
    autor_p = autor_frame.paragraphs[0]
    autor_p.alignment = PP_ALIGN.CENTER
    autor_p.font.size = Pt(16)
    autor_p.font.color.rgb = RGBColor(200, 200, 200)

def adicionar_slide_problema(slide):
    """Slide do problema"""
    adicionar_titulo(slide, "O Problema")
    
    conteudo = """🏪 Loja Online com Tráfego Variável

📊 Padrão Observado:
• Segunda a Quinta: Tráfego normal (~1.000 usuários/hora)
• Sexta-feira 18h-22h: PICO (10.000+ usuários/hora) 
• Fim de semana: Baixíssimo tráfego (~200 usuários/hora)

⚠️ Problemas:
• Auto Scaling SEM LIMITES → Custos descontrolados
• Servidores ociosos no fim de semana → Desperdício
• Infraestrutura manual → Lento para ajustar
• Sem controle de versão → Risco operacional
"""
    adicionar_texto_box(slide, conteudo, Inches(1), Inches(1.5), Inches(8), Inches(5))

def adicionar_slide_pico_trafego(slide):
    """Slide do pico de tráfego"""
    adicionar_titulo(slide, "Sexta-feira 20h 🔥")
    
    conteudo = """
📈 PICO DE TRÁFEGO

Sistema antes do Terraform:
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

2 servidores → Sistema LENTO
Auto Scaling dispara...
→ 5 servidores
→ 10 servidores  
→ 20 servidores
→ 35 servidores 😱

❌ Problema: Sem limite máximo definido!
"""
    adicionar_texto_box(slide, conteudo, Inches(1.5), Inches(1.5), Inches(7), Inches(5))

def adicionar_slide_custo(slide):
    """Slide do custo"""
    adicionar_titulo(slide, "Segunda-feira 8h 💸")
    
    conteudo = """
CEO abre o email da AWS...

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Fatura AWS - Outubro 2024
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

35 instâncias × 72 horas × $0.096/hora = $242
(Fim de semana inteiro com servidores ociosos!)

10 instâncias rodando 24/7 quando precisava 2

Total mensal: ~$50.000 💰

❌ INSUSTENTÁVEL!
"""
    texto_box = adicionar_texto_box(slide, conteudo, Inches(1), Inches(1.5), Inches(8), Inches(5))
    # Última linha em vermelho
    paragrafos = texto_box.text_frame.paragraphs
    paragrafos[-1].font.color.rgb = ERROR_RED
    paragrafos[-1].font.size = Pt(32)
    paragrafos[-1].font.bold = True

def adicionar_slide_diagrama(slide, titulo, caminho_imagem):
    """Slide com diagrama"""
    adicionar_titulo(slide, titulo)
    try:
        slide.shapes.add_picture(caminho_imagem, Inches(0.5), Inches(1.5), 
                                width=Inches(9), height=Inches(5.5))
    except:
        # Se não encontrar imagem, mostra placeholder
        texto = f"[Diagrama: {titulo}]\n\nAbra o arquivo:\n{caminho_imagem}"
        adicionar_texto_box(slide, texto, Inches(2), Inches(2.5), Inches(6), Inches(3))

def adicionar_slide_estrutura(slide):
    """Slide da estrutura do projeto"""
    adicionar_titulo(slide, "Estrutura do Projeto Terraform")
    
    conteudo = """terraform/localstack/
├── provider.tf      ← Configuração LocalStack
├── variables.tf     ← CONTROLE DE CUSTOS aqui! 💰
├── vpc.tf           ← Rede (VPC, Subnets, NAT)
├── ec2.tf           ← Compute (ALB, ASG, Policies)
├── data.tf          ← Data sources (AMI, etc)
└── outputs.tf       ← Informações úteis

📝 Cada arquivo tem UMA responsabilidade
✅ Código versionado no Git
✅ Reproduzível em qualquer ambiente
"""
    adicionar_codigo(slide, conteudo, Inches(1.5), Inches(1.5), Inches(7), Inches(5))

def adicionar_slide_variables(slide):
    """Slide de variáveis - controle de custos"""
    adicionar_titulo(slide, "variables.tf - O Controle de Custos 💰")
    
    codigo = """# Auto Scaling - LIMITES DEFINIDOS
variable "min_size" {
  default = 2        # Mínimo para alta disponibilidade
}

variable "max_size" {
  default = 10       # ⚠️ MÁXIMO = Proteção de custos!
}

variable "desired_capacity" {
  default = 2        # Capacidade normal
}

# Scheduled Scaling - Padrões Previsíveis
variable "friday_desired_capacity" {
  default = 6        # Sexta-feira: aumenta
}

variable "weekend_desired_capacity" {
  default = 2        # Fim de semana: reduz
}
"""
    adicionar_codigo(slide, codigo, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))

def adicionar_slide_vpc(slide):
    """Slide da VPC"""
    adicionar_titulo(slide, "vpc.tf - A Base da Rede 🏗️")
    
    codigo = """# VPC = Virtual Private Cloud
# "Nosso datacenter virtual na AWS"

resource "aws_vpc" "main" {
  cidr_block           = "10.0.0.0/16"     # 65.536 IPs
  enable_dns_hostnames = true
  enable_dns_support   = true
  
  tags = {
    Name = "loja-online-prod-vpc"
  }
}

# Internet Gateway = "Porta de entrada"
resource "aws_internet_gateway" "main" {
  vpc_id = aws_vpc.main.id
  
  tags = {
    Name = "loja-online-prod-igw"
  }
}
"""
    adicionar_codigo(slide, codigo, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))

def adicionar_slide_subnets(slide):
    """Slide de subnets"""
    adicionar_titulo(slide, "Subnets - As 3 Camadas 🏢")
    
    conteudo = """
┌─────────────────────────────────────────────┐
│ 🌐 Subnets PÚBLICAS (10.0.1-3.0/24)        │
│    → Load Balancer, NAT Gateway            │
│    → Tem acesso à Internet                 │
└─────────────────────────────────────────────┘
         ↓
┌─────────────────────────────────────────────┐
│ 🔒 Subnets PRIVADAS (10.0.11-13.0/24)     │
│    → Servidores da aplicação (EC2)         │
│    → Acesso via NAT Gateway                │
└─────────────────────────────────────────────┘
         ↓
┌─────────────────────────────────────────────┐
│ 🔐 Subnets DATABASE (10.0.21-23.0/24)     │
│    → Banco de dados (RDS)                  │
│    → SEM acesso à Internet                 │
│    → Máxima segurança                      │
└─────────────────────────────────────────────┘

✅ 3 Availability Zones = Alta Disponibilidade
"""
    adicionar_codigo(slide, conteudo, Inches(1), Inches(1.5), Inches(8), Inches(5.5))

def adicionar_slide_alb(slide):
    """Slide do Load Balancer"""
    adicionar_titulo(slide, "Application Load Balancer - O Porteiro 🚪")
    
    codigo = """# ALB = Distribui tráfego entre servidores

resource "aws_lb" "main" {
  name               = "loja-online-prod-alb"
  internal           = false           # Público
  load_balancer_type = "application"
  
  subnets = [
    aws_subnet.public_1a.id,
    aws_subnet.public_1b.id,
    aws_subnet.public_1c.id
  ]
}

# Target Group = "Para onde enviar as requisições"
resource "aws_lb_target_group" "backend" {
  name     = "loja-online-backend-tg"
  port     = 3000
  protocol = "HTTP"
  vpc_id   = aws_vpc.main.id
  
  health_check {
    path = "/health"    # Verifica se servidor está OK
  }
}
"""
    adicionar_codigo(slide, codigo, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))

def adicionar_slide_asg(slide):
    """Slide do Auto Scaling Group"""
    adicionar_titulo(slide, "Auto Scaling Group - A Mágica ✨")
    
    codigo = """# ASG = Ajusta número de servidores automaticamente

resource "aws_autoscaling_group" "backend" {
  name                = "loja-online-backend-asg"
  vpc_zone_identifier = [/* subnets privadas */]
  
  min_size         = var.min_size          # 2
  max_size         = var.max_size          # 10 ⚠️
  desired_capacity = var.desired_capacity  # 2
  
  health_check_type         = "ELB"
  health_check_grace_period = 300
  
  target_group_arns = [aws_lb_target_group.backend.arn]
  
  launch_template {
    id      = aws_launch_template.backend.id
    version = "$Latest"
  }
}

# 🎯 Target Tracking = Escala baseado em CPU
resource "aws_autoscaling_policy" "cpu_tracking" {
  policy_type            = "TargetTrackingScaling"
  target_value           = 70.0   # CPU > 70% → Adiciona servidor
}
"""
    adicionar_codigo(slide, codigo, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))

def adicionar_slide_scheduled(slide):
    """Slide de Scheduled Scaling"""
    adicionar_titulo(slide, "Scheduled Scaling - Horários Especiais ⏰")
    
    codigo = """# Sexta-feira 18h - AUMENTA para o pico
resource "aws_autoscaling_schedule" "friday_scale_up" {
  scheduled_action_name  = "friday-evening-scale-up"
  autoscaling_group_name = aws_autoscaling_group.backend.name
  
  recurrence       = "0 18 * * 5"    # Cron: Sexta 18h
  desired_capacity = 6                # 2 → 6 servidores
}

# Sexta-feira 23h - REDUZ após o pico
resource "aws_autoscaling_schedule" "friday_scale_down" {
  scheduled_action_name  = "friday-night-scale-down"
  autoscaling_group_name = aws_autoscaling_group.backend.name
  
  recurrence       = "0 23 * * 5"    # Cron: Sexta 23h
  desired_capacity = 2                # 6 → 2 servidores
}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Timeline:
Segunda-Quinta: 2 servidores (normal)
Sexta 18h:      6 servidores (pico esperado)
Sexta 23h:      2 servidores (volta ao normal)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
"""
    adicionar_codigo(slide, codigo, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))

def adicionar_slide_security(slide):
    """Slide de Security Groups"""
    adicionar_titulo(slide, "Security Groups - O Firewall 🔒")
    
    codigo = """# Security Group = Firewall de cada recurso

# SG do Load Balancer - Aceita internet
resource "aws_security_group" "alb" {
  name        = "loja-online-alb-sg"
  description = "ALB security group"
  vpc_id      = aws_vpc.main.id
  
  ingress {
    from_port   = 80     # HTTP
    to_port     = 80
    protocol    = "tcp"
    cidr_blocks = ["0.0.0.0/0"]  # Qualquer IP
  }
  
  ingress {
    from_port   = 443    # HTTPS
    to_port     = 443
    protocol    = "tcp"
    cidr_blocks = ["0.0.0.0/0"]
  }
}

# SG dos Servidores - APENAS do ALB
resource "aws_security_group" "ec2" {
  name = "loja-online-ec2-sg"
  
  ingress {
    from_port       = 3000
    to_port         = 3000
    protocol        = "tcp"
    security_groups = [aws_security_group.alb.id]  # ⚠️ Só ALB!
  }
}
"""
    adicionar_codigo(slide, codigo, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))

def adicionar_slide_demo(slide):
    """Slide de demo"""
    adicionar_titulo(slide, "Demo - Vamos Executar! 🚀")
    
    conteudo = """
1️⃣ Iniciar LocalStack
   $ docker compose up -d

2️⃣ Inicializar Terraform
   $ cd terraform/localstack
   $ terraform init

3️⃣ Ver o que será criado
   $ terraform plan

4️⃣ Aplicar a infraestrutura
   $ terraform apply
   → Digite: yes

5️⃣ Verificar recursos criados
   $ terraform output
   $ terraform state list

6️⃣ Destruir (ao final)
   $ terraform destroy

⏱️ Tempo total: ~3 minutos
✅ 20+ recursos criados automaticamente!
"""
    adicionar_codigo(slide, conteudo, Inches(1.5), Inches(1.5), Inches(7), Inches(5.5))

def adicionar_slide_comparacao(slide):
    """Slide de comparação"""
    adicionar_titulo(slide, "Antes vs Depois")
    
    # Coluna ANTES (vermelho)
    antes = """❌ ANTES

• 10 servidores 24/7
• R$ 50.000/mês
• Infraestrutura manual
• Sem versionamento
• Quedas nas sextas
• Desperdício fim de semana
• Deployment manual
• Rollback difícil
"""
    antes_box = adicionar_texto_box(slide, antes, Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    for p in antes_box.text_frame.paragraphs:
        p.font.color.rgb = ERROR_RED
    
    # Coluna DEPOIS (verde)
    depois = """✅ DEPOIS

• 2-10 servidores
• R$ 18.000/mês (64% ↓)
• Código Terraform
• Git versionado
• Escala automaticamente
• Scheduled scaling
• terraform apply
• Fácil rollback
"""
    depois_box = adicionar_texto_box(slide, depois, Inches(5.5), Inches(1.5), Inches(4), Inches(5.5))
    for p in depois_box.text_frame.paragraphs:
        p.font.color.rgb = SUCCESS_GREEN

def adicionar_slide_resultados(slide):
    """Slide de resultados"""
    adicionar_titulo(slide, "Resultados Alcançados 🎯")
    
    conteudo = """
💰 ECONOMIA DE CUSTOS
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
   Antes: R$ 50.000/mês
   Depois: R$ 18.000/mês
   Economia: R$ 32.000/mês (64%)


⚡ PERFORMANCE
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
   ✅ Sextas-feiras: Sistema estável
   ✅ Auto Scaling controlado (max: 10)
   ✅ Health checks automáticos


🔒 SEGURANÇA
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
   ✅ 3 camadas de rede isoladas
   ✅ Security Groups restritivos
   ✅ IAM com menor privilégio


🚀 OPERACIONAL
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
   ✅ Deploy em 3 minutos
   ✅ Código versionado no Git
   ✅ Reproduzível em qualquer ambiente
"""
    adicionar_texto_box(slide, conteudo, Inches(1), Inches(1.5), Inches(8), Inches(5.5))

def adicionar_slide_frontend_codigo(slide):
    """Slide de código do Frontend"""
    adicionar_titulo(slide, "Frontend - CloudFront + ALB 🌐")
    
    codigo = """# CloudFront CDN - Distribuição Global (comentado no LocalStack)
resource "aws_cloudfront_distribution" "main" {
  enabled             = true
  comment             = "Loja Online CDN"
  default_root_object = "index.html"
  
  origin {
    domain_name = aws_lb.main.dns_name    # ALB como origem
    origin_id   = "alb-origin"
  }
  
  # Cache otimizado
  default_cache_behavior {
    allowed_methods        = ["GET", "HEAD", "OPTIONS"]
    cached_methods         = ["GET", "HEAD"]
    target_origin_id       = "alb-origin"
    viewer_protocol_policy = "redirect-to-https"
    
    min_ttl     = 0
    default_ttl = 3600      # 1 hora
    max_ttl     = 86400     # 24 horas
  }
}

# Route53 - DNS
resource "aws_route53_zone" "main" {
  name = "loja-online.com.br"
}

# WAF - Web Application Firewall (comentado)
resource "aws_wafv2_web_acl" "main" {
  name  = "loja-online-waf"
  scope = "CLOUDFRONT"
  
  # Regras de proteção
}
"""
    adicionar_codigo(slide, codigo, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))

def adicionar_slide_backend_codigo(slide):
    """Slide de código do Backend"""
    adicionar_titulo(slide, "Backend - Aplicação Node.js 🚀")
    
    codigo = """# Launch Template - Configuração dos Servidores
resource "aws_launch_template" "backend" {
  name_prefix   = "loja-online-backend-"
  image_id      = data.aws_ami.amazon_linux_2023.id
  instance_type = "t3.micro"
  
  # User Data - Script de inicialização
  user_data = base64encode(<<-EOF
    #!/bin/bash
    # Instalar Node.js 20
    curl -fsSL https://rpm.nodesource.com/setup_20.x | bash -
    yum install -y nodejs
    
    # Aplicação Express
    cat > /home/ec2-user/app.js << 'APP'
    const express = require('express');
    const app = express();
    
    app.get('/health', (req, res) => {
      res.json({ status: 'healthy' });
    });
    
    app.listen(3000);
    APP
    
    # Iniciar aplicação
    npm install express
    node /home/ec2-user/app.js
  EOF
  )
  
  # IAM Role
  iam_instance_profile {
    name = aws_iam_instance_profile.ec2.name
  }
}

# SQS - Fila de Mensagens (comentado)
resource "aws_sqs_queue" "orders" {
  name = "loja-online-orders-queue"
}
"""
    adicionar_codigo(slide, codigo, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))

def adicionar_slide_database_codigo(slide):
    """Slide de código do Database"""
    adicionar_titulo(slide, "Database - RDS + ElastiCache 💾")
    
    codigo = """# RDS PostgreSQL Multi-AZ (comentado no LocalStack)
resource "aws_db_instance" "main" {
  identifier     = "loja-online-db"
  engine         = "postgres"
  engine_version = "15.4"
  instance_class = "db.t3.micro"
  
  allocated_storage     = 20
  storage_encrypted     = true      # Criptografia
  
  multi_az              = true      # Alta Disponibilidade
  db_subnet_group_name  = aws_db_subnet_group.main.name
  
  backup_retention_period = 7       # 7 dias de backup
  backup_window          = "03:00-04:00"
  maintenance_window     = "sun:04:00-sun:05:00"
  
  vpc_security_group_ids = [aws_security_group.database.id]
}

# Read Replica - Leitura
resource "aws_db_instance" "read_replica" {
  replicate_source_db = aws_db_instance.main.identifier
  instance_class      = "db.t3.micro"
  
  availability_zone = "us-east-1b"  # AZ diferente
}

# ElastiCache Redis - Cache (comentado)
resource "aws_elasticache_cluster" "redis" {
  cluster_id           = "loja-online-cache"
  engine               = "redis"
  node_type            = "cache.t3.micro"
  num_cache_nodes      = 1
  parameter_group_name = "default.redis7"
  port                 = 6379
}

# Subnet Group - Database isolado
resource "aws_db_subnet_group" "main" {
  name       = "loja-online-db-subnet-group"
  subnet_ids = [
    aws_subnet.database_1a.id,
    aws_subnet.database_1b.id,
    aws_subnet.database_1c.id
  ]
}
"""
    adicionar_codigo(slide, codigo, Inches(0.4), Inches(1.5), Inches(9.2), Inches(5.5))

def adicionar_slide_proximos_passos(slide):
    """Slide de próximos passos"""
    adicionar_titulo(slide, "Próximos Passos 📚")
    
    conteudo = """
🎯 PARA PRATICAR:

1. Clone o repositório
   git clone <repo>

2. Execute o projeto com LocalStack
   Sem custos! Tudo local.

3. Modifique as variáveis
   Experimente: min_size, max_size, etc.

4. Adicione novos recursos
   RDS, ElastiCache, CloudFront...


📖 PARA APROFUNDAR:

• Terraform Registry: registry.terraform.io
• AWS Well-Architected: aws.amazon.com/architecture
• Terraform Best Practices
• AWS Cost Optimization
• LocalStack Documentation
"""
    adicionar_texto_box(slide, conteudo, Inches(1), Inches(1.5), Inches(8), Inches(5.5))

# Funções auxiliares
def adicionar_titulo(slide, texto):
    """Adiciona título ao slide"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = texto
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = AWS_BLUE
    
    # Linha decorativa
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = AWS_ORANGE
    line.line.fill.background()

def adicionar_texto_box(slide, texto, left, top, width, height):
    """Adiciona caixa de texto"""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = texto
    text_frame.word_wrap = True
    
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.name = 'Consolas'
    
    return text_box

def adicionar_codigo(slide, codigo, left, top, width, height):
    """Adiciona bloco de código"""
    code_box = slide.shapes.add_textbox(left, top, width, height)
    code_frame = code_box.text_frame
    code_frame.text = codigo
    code_frame.word_wrap = False
    
    for paragraph in code_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.name = 'Consolas'
        paragraph.font.color.rgb = RGBColor(40, 40, 40)
    
    # Fundo cinza claro
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    return code_box

if __name__ == "__main__":
    criar_apresentacao()
